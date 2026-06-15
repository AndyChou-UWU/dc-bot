"""
生成專業的 Discord AI Bot 項目 PPT
使用 python-pptx 庫
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# 創建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 定義顏色方案
COLOR_DARK = RGBColor(25, 25, 112)  # 午夜藍
COLOR_ACCENT = RGBColor(65, 105, 225)  # 皇家藍
COLOR_HIGHLIGHT = RGBColor(255, 69, 0)  # 橙紅
COLOR_TEXT = RGBColor(50, 50, 50)
COLOR_LIGHT_BG = RGBColor(240, 248, 255)

def add_title_slide(prs, title, subtitle=""):
    """添加標題頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 副標題
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = COLOR_HIGHLIGHT
    
    return slide

def add_content_slide(prs, title, content_items):
    """添加內容頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題欄
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_ACCENT
    title_shape.line.color.rgb = COLOR_ACCENT
    
    # 標題文字
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # 內容
    left = Inches(0.8)
    top = Inches(1.5)
    for idx, item in enumerate(content_items):
        text_box = slide.shapes.add_textbox(left, top + idx * Inches(0.9), Inches(8.4), Inches(0.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        if isinstance(item, tuple):
            title_text, desc = item
            text_frame.text = f"{title_text}"
            p = text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = COLOR_DARK
            
            # 描述
            p2 = text_frame.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(16)
            p2.font.color.rgb = COLOR_TEXT
            p2.level = 0
        else:
            text_frame.text = f"• {item}"
            p = text_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_TEXT
            p.space_before = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """添加兩列內容頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_ACCENT
    title_shape.line.color.rgb = COLOR_ACCENT
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # 左列
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for idx, item in enumerate(left_items):
        if idx == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = f"• {item}" if idx > 0 else item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
    
    # 右列
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for idx, item in enumerate(right_items):
        if idx == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = f"• {item}" if idx > 0 else item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
    
    return slide

# ===================== 開始創建 PPT =====================

# 第1頁：標題頁
add_title_slide(prs, "🤖 Discord AI 聊天機器人", 
                "Python 大作業 | 多功能 AI 助手系統")

# 第2頁：項目概述
add_content_slide(prs, "📌 項目概述", [
    ("項目名稱", "subaso-俗北ㄙㄡˊ Discord AI Bot"),
    ("主要目標", "爲 Discord 用戶提供智能化的 AI 聊天和娛樂功能"),
    ("核心技術", "Python 3.8+ | discord.py | Ollama 本地 AI | 異步編程"),
    ("版本", "v2.0.0"),
])

# 第3頁：核心功能
add_content_slide(prs, "⚡ 核心功能特色", [
    "🎭 5 種 AI 角色 - 閒談、數理、語文、程式、家務",
    "💬 連續對話 - 記住對話歷史，支持多輪交互",
    "🌍 多語言支持 - 繁體中文、英文、日文、韓文、西班牙文",
    "📝 數據持久化 - 自動保存用戶偏好和對話記錄",
    "⚡ 本地 AI 引擎 - 基於 Ollama，無需付費 API，隱私有保障",
    "🔔 版本管理 - 實時推送更新通知給所有用戶",
])

# 第4頁：功能模塊
add_content_slide(prs, "🎯 功能模塊架構", [
    ("🤖 AI 對話系統", "多角色、多語言、連續對話、版本管理、數據保存"),
    ("💰 經濟系統", "挖礦、釣魚、狩獵、賭博、寵物系統、市場交易"),
    ("🎮 娛樂遊戲", "Pokemon 猜謎、Trivia 問答、動漫角色、數字遊戲"),
    ("⚙️ 管理工具", "日誌管理、消息清理、歡迎消息、自定義命令"),
])

# 第5頁：AI 角色詳解
add_two_column_slide(prs, "🎭 5 大 AI 角色詳解",
    ["✅ 閒談 - 友善幽默的日常聊天夥伴",
     "🔢 數理 - 數學和科學問題專家",
     "📚 語文 - 語言文學和寫作指導"],
    ["💻 程式 - 編程技術問題解決",
     "🏠 家務 - 生活烹飪清潔建議"])

# 第6頁：技術架構
add_content_slide(prs, "🔧 技術棧與架構", [
    ("編程語言", "Python 3.8+ | 異步編程 (asyncio)"),
    ("主要庫", "discord.py 2.0+ | aiohttp | httpx | python-dotenv"),
    ("AI 引擎", "Ollama (本地輕量級 LLM) | Qwen2.5-1.5B 模型"),
    ("數據存儲", "JSON 文件存儲 | user_data.json 用戶數據"),
    ("日誌系統", "結構化日誌 | 文件和控制檯輸出 | UTF-8 編碼支持"),
])

# 第7頁：系統架構圖（文字描述）
add_two_column_slide(prs, "📐 系統架構",
    ["Discord API",
     "↓",
     "bot.py (主程序)",
     "↓",
     "prism_config.py",
     "↓",
     "Ollama AI"],
    ["用戶交互",
     "↓",
     "命令處理",
     "↓",
     "AI 角色切換",
     "↓",
     "本地智能回覆"])

# 第8頁：工作流程
add_content_slide(prs, "🔄 工作流程", [
    "1️⃣ 用戶在 Discord 發送消息或命令",
    "2️⃣ Bot 接收事件，解析命令和參數",
    "3️⃣ 根據選定的 AI 角色，構建 system prompt",
    "4️⃣ 調用 Ollama API，發送消息和對話歷史",
    "5️⃣ Ollama 本地推理，返回 AI 回覆",
    "6️⃣ Bot 格式化回覆，發送到 Discord",
    "7️⃣ 保存對話記錄和用戶數據到本地文件",
])

# 第9頁：使用示例
add_content_slide(prs, "💡 使用示例", [
    ("數學問題", "/ai 數理 如何求解二次方程？"),
    ("日常閒聊", "/ai 閒談 今天天氣真好！"),
    ("編程幫助", "/ai 程式 Python 中如何定義類？"),
    ("語文指導", "/ai 語文 怎樣寫出好的開頭？"),
    ("生活建議", "/ai 家務 如何快速清潔廚房？"),
])

# 第10頁：技術亮點
add_content_slide(prs, "✨ 技術亮點", [
    "🔐 異步編程 - 高效處理多個併發用戶請求",
    "🎯 面向對象設計 - 清晰的模塊劃分和可擴展性",
    "🔄 狀態管理 - 完整的用戶數據和對話歷史管理",
    "🌐 API 集成 - 與 Discord 和 Ollama 的無縫對接",
    "📊 日誌系統 - 結構化日誌便於調試和監控",
    "🔧 配置管理 - 集中配置，易於自定義和維護",
])

# 第11頁：環境要求
add_two_column_slide(prs, "📋 環境與依賴",
    ["硬件要求:",
     "• Python 3.8+",
     "• 8GB+ 內存",
     "• 15GB+ 磁盤空間",
     "• 網絡連接"],
    ["軟件依賴:",
     "• discord.py >= 2.0.0",
     "• python-dotenv",
     "• aiohttp",
     "• Ollama (本地 AI)",
     "• Discord Token"])

# 第12頁：部署方式
add_content_slide(prs, "🚀 部署與運行", [
    "1️⃣ 安裝依賴：pip install -r requirements.txt",
    "2️⃣ 安裝 Ollama：https://ollama.ai",
    "3️⃣ 拉取模型：ollama pull Qwen2.5-1.5B",
    "4️⃣ 配置 Token：在 設定.env 中設置 DISCORD_TOKEN",
    "5️⃣ 啓動 Bot：python bot.py",
    "6️⃣ Docker 部署：docker build -t dcard-bot . && docker run ...",
])

# 第13頁：改進與展望
add_content_slide(prs, "🎯 改進與展望", [
    "🔮 未來計劃",
    "• 支持更多 AI 模型（GPT、Claude 等付費 API）",
    "• 添加語音交互功能",
    "• 實現用戶權限管理系統",
    "• 開發 Web 控制面板",
    "• 支持數據庫存儲（MongoDB、PostgreSQL）",
    "• 多服務器部署和負載均衡",
])

# 第14頁：項目成果
add_content_slide(prs, "🏆 項目成果與創新", [
    ("完整性", "完全自主開發，功能模塊齊全，代碼規範清晰"),
    ("創新性", "結合本地 AI 和 Discord Bot，提供隱私保護的智能助手"),
    ("可用性", "開箱即用，配置簡單，文檔詳細"),
    ("可擴展性", "模塊化設計，易於添加新功能和新角色"),
    ("學習價值", "綜合運用 Python 異步編程、API 開發、數據管理等技能"),
])

# 第15頁：結論
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_DARK

conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3.5))
conclusion_frame = conclusion_box.text_frame
conclusion_frame.word_wrap = True

p1 = conclusion_frame.paragraphs[0]
p1.text = "感謝評審！"
p1.font.size = Pt(48)
p1.font.bold = True
p1.font.color.rgb = COLOR_HIGHLIGHT
p1.alignment = PP_ALIGN.CENTER

p2 = conclusion_frame.add_paragraph()
p2.text = "\n本項目展示了現代 Python 開發的最佳實踐\n以及 AI 應用的實際部署能力"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(255, 255, 255)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

p3 = conclusion_frame.add_paragraph()
p3.text = "\n源代碼已上傳 GitHub"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(200, 200, 200)
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(40)

# 保存 PPT
output_path = "Discord_AI_Bot_大作業_演示.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"📊 共 {len(prs.slides)} 頁幻燈片")
print(f"⏰ 生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
