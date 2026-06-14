"""
生成小學生版 PowerPoint - 用簡單比喻解釋程式邏輯
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ 缺少 python-pptx，正在安裝...")
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

PROJECT_DIR = Path(__file__).parent
PPT_OUTPUT = PROJECT_DIR / "小學生版_程式邏輯說明.pptx"

# 彩色配置
COLORS = {
    'title': RGBColor(52, 73, 94),
    'accent1': RGBColor(231, 76, 60),      # 紅
    'accent2': RGBColor(52, 152, 219),    # 藍
    'accent3': RGBColor(46, 204, 113),    # 綠
    'accent4': RGBColor(241, 196, 15),    # 黃
    'accent5': RGBColor(155, 89, 182),    # 紫
    'text': RGBColor(44, 62, 80),
    'light_bg': RGBColor(236, 240, 241)
}

def add_title_slide(prs, title, subtitle, emoji=""):
    """簡單標題幻燈片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['accent2']
    
    # 大標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = emoji + "\n" + title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    for line in subtitle.split('\n'):
        p = subtitle_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, emoji, content_list):
    """內容幻燈片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題背景
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['accent2']
    title_shape.line.color.rgb = COLORS['accent2']
    
    # 標題文字
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = emoji + "  " + title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def add_comparison_slide(prs, title, emoji, left_title, left_items, right_title, right_items, left_color, right_color):
    """對比幻燈片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['accent2']
    title_shape.line.color.rgb = COLORS['accent2']
    
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = emoji + "  " + title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 左邊盒子
    left_box_shape = slide.shapes.add_shape(1, Inches(0.3), Inches(1.2), Inches(4.5), Inches(5.8))
    left_box_shape.fill.solid()
    left_box_shape.fill.fore_color.rgb = left_color
    left_box_shape.line.color.rgb = left_color
    left_box_shape.line.width = Pt(2)
    
    left_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.1), Inches(5.4))
    left_frame = left_text.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0
    
    # 右邊盒子
    right_box_shape = slide.shapes.add_shape(1, Inches(5.2), Inches(1.2), Inches(4.5), Inches(5.8))
    right_box_shape.fill.solid()
    right_box_shape.fill.fore_color.rgb = right_color
    right_box_shape.line.color.rgb = right_color
    right_box_shape.line.width = Pt(2)
    
    right_text = slide.shapes.add_textbox(Inches(5.4), Inches(1.4), Inches(4.1), Inches(5.4))
    right_frame = right_text.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0

def create_elementary_ppt():
    """創建小學生版 PowerPoint"""
    print("🎬 正在生成小學生版 PowerPoint...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==================== 第 1 張：封面 ====================
    add_title_slide(prs,
        "程式邏輯\n小學生版",
        "用簡單比喻理解\nBot 是怎麼工作的",
        "🤖")
    
    # ==================== 第 2 張：什麼是 Bot ====================
    add_content_slide(prs, "什麼是 Bot?", "🤖", [
        "🎯 Bot 是一個「虛擬助手」",
        "",
        "✓ 能和你聊天",
        "✓ 能玩遊戲",
        "✓ 能記住你的喜好",
        "✓ 能管理虛擬金幣",
        "",
        "就像你的朋友一樣！",
        "只是它住在電腦裡，不是真人"
    ])
    
    # ==================== 第 3 張：Bot 的大腦 ====================
    add_content_slide(prs, "Bot 的「大腦」", "🧠", [
        "你說話 ➜ Bot 聽到 ➜ Bot 思考 ➜ Bot 回答",
        "",
        "就像這個流程：",
        "",
        "1️⃣ 你：!mode 閒談",
        "2️⃣ Bot 理解：要切換角色",
        "3️⃣ Bot 記下來：這個人用閒談模式",
        "4️⃣ Bot 回覆：✅ 已切換到閒談模式",
        "",
        "然後下一次你提問時，",
        "Bot 就會用「閒談」的方式回答"
    ])
    
    # ==================== 第 4 張：五個人格 ====================
    add_content_slide(prs, "五個不同的人格", "🎭", [
        "✅ 閒談 - 愛聊天，很友善",
        "",
        "🔢 數理 - 會計算，很聰明",
        "",
        "📚 語文 - 愛讀書，很文化",
        "",
        "💻 程式 - 很厲害，會修東西",
        "",
        "🏠 家務 - 很會做飯和打掃",
        "",
        "就像你有五個朋友，性格都不一樣！"
    ])
    
    # ==================== 第 5 張：對比 - 沒記憶 vs 有記憶 ====================
    add_comparison_slide(prs, "為什麼 Bot 要有記憶?", "💾",
        "沒有記憶", [
            "你：我叫小明",
            "Bot：很高興認識",
            "",
            "你：我喜歡吃冰",
            "Bot：冰是什麼？",
            "",
            "❌ 每次都要重新說"
        ], 
        "有記憶", [
            "你：我叫小明",
            "Bot：很高興認識",
            "",
            "你：我喜歡吃冰",
            "Bot：小明，我記得你",
            "",
            "✅ 能記住你的事"
        ],
        COLORS['accent1'], COLORS['accent3'])
    
    # ==================== 第 6 張：虛擬金幣遊戲 ====================
    add_content_slide(prs, "虛擬金幣系統", "💰", [
        "就像手機遊戲有虛擬金幣一樣",
        "",
        "⛏️ 挖礦 → 得到 50-500 金幣",
        "🐟 釣魚 → 得到 40-400 金幣",
        "🎰 賭博 → 可能贏或輸",
        "🐾 寵物 → 需要照顧",
        "",
        "開始時有 1000 金幣",
        "可以用來玩遊戲、養寵物",
        "",
        "⚠️ 注意：有冷卻時間（等一下再挖）",
        "這樣才公平，誰都不能無限賺錢！"
    ])
    
    # ==================== 第 7 張：遊戲系統 ====================
    add_content_slide(prs, "遊戲有多好玩？", "🎮", [
        "🎯 Pokemon 猜謎",
        "   我會噴火 → 答：火焰鼠 ✓ 加 100 分",
        "",
        "🧠 Trivia 知識賽",
        "   答對問題就得分",
        "",
        "🎬 動漫角色猜謎",
        "   根據描述猜人物",
        "",
        "🔢 數字遊戲",
        "   猜我想的數字",
        "",
        "✨ 互動動作",
        "   !hug !pat !dance"
    ])
    
    # ==================== 第 8 張：指令就是魔法咒語 ====================
    add_content_slide(prs, "指令 = 魔法咒語", "✨", [
        "Bot 只聽懂特定的指令",
        "",
        "!mode 閒談     → 切換角色",
        "!ask 你好      → 問問題",
        "!balance       → 看金幣",
        "!mine          → 挖礦",
        "!help          → 看所有指令",
        "",
        "就像魔法咒語一樣：",
        "說對了才有效，說錯就沒反應"
    ])
    
    # ==================== 第 9 張：存檔的重要性 ====================
    add_comparison_slide(prs, "為什麼要存檔？", "💾",
        "沒有存檔", [
            "遊戲重開",
            "",
            "❌ 金幣消失",
            "❌ 角色重設",
            "❌ 寵物沒了",
            "",
            "很麻煩"
        ], 
        "有存檔", [
            "遊戲重開",
            "",
            "✅ 金幣保存",
            "✅ 角色記住",
            "✅ 寵物還在",
            "",
            "很方便"
        ],
        COLORS['accent1'], COLORS['accent3'])
    
    # ==================== 第 10 張：三個重要概念 ====================
    add_content_slide(prs, "程式的三個重要概念", "🎓", [
        "1️⃣ 順序 - 按先後順序做事",
        "   起床 → 洗臉 → 吃飯 → 上學",
        "   順序不能亂，不然會出問題",
        "",
        "2️⃣ 判斷 - 根據情況做不同的事",
        "   如果下雨 → 帶傘",
        "   如果晴天 → 帶眼鏡",
        "",
        "3️⃣ 重複 - 做很多次相同的事",
        "   重複 100 次寫生字",
        "   Bot 重複接訊息和回覆"
    ])
    
    # ==================== 第 11 張：Bot 啟動過程 ====================
    add_content_slide(prs, "Bot 啟動流程", "🚀", [
        "python bot.py",
        "    ↓",
        "Bot 讀取設定",
        "    ↓",
        "Bot 連接到 Discord",
        "    ↓",
        "Bot 上線！等待訊息",
        "    ↓",
        "接到訊息 → 理解 → 回覆",
        "    ↓",
        "重複等待和回覆",
        "",
        "按 Ctrl+C 才會停止"
    ])
    
    # ==================== 第 12 張：程式 vs 生活 ====================
    add_comparison_slide(prs, "程式和生活的相似性", "🌍",
        "現實生活", [
            "👨‍🏫 老師教你做事",
            "",
            "📖 你要記住規則",
            "",
            "🤔 下雨你就帶傘",
            "",
            "📝 重複練習",
            "",
            "💾 記得朋友的名字"
        ], 
        "程式邏輯", [
            "👨‍💻 程式員寫指令",
            "",
            "🤖 Bot 要遵守規則",
            "",
            "⚡ 訊息是 !ask 就回答",
            "",
            "🔄 重複接收和回覆",
            "",
            "🗂️ 記得使用者的資料"
        ],
        COLORS['accent4'], COLORS['accent5'])
    
    # ==================== 第 13 張：不要怕，其實很簡單 ====================
    add_content_slide(prs, "其實程式不難！", "💡", [
        "程式看起來很複雜，但其實就是：",
        "",
        "✅ 告訴機器做什麼（指令）",
        "✅ 記住重要資訊（記憶）",
        "✅ 根據情況做不同的事（判斷）",
        "✅ 重複做同樣的事（循環）",
        "",
        "你每天都在做這些事：",
        "· 按順序上學",
        "· 根據天氣穿衣",
        "· 重複寫功課",
        "· 記得媽媽的話",
        "",
        "程式只是用更嚴格、更精確的方式"
    ])
    
    # ==================== 第 14 張：最終結論 ====================
    add_title_slide(prs,
        "編程 = 教機器做事",
        "就像教弟弟妹妹一樣！\n\n簡單、有趣、又有用",
        "🎉")
    
    prs.save(str(PPT_OUTPUT))
    print(f"✅ 小學生版 PowerPoint 已生成: {PPT_OUTPUT}")
    print(f"📊 共 {len(prs.slides)} 張幻燈片")
    return True

if __name__ == "__main__":
    try:
        create_elementary_ppt()
        print("\n🎉 完成！")
    except Exception as e:
        print(f"❌ 生成失敗: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
