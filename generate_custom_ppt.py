"""
自訂 PPT 產生腳本 - 針對頻審的功能介紹、技術架構、使用指南
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    # 若未安裝 python-pptx，則自動安裝
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

# 專案根目錄
PROJECT_DIR = Path(__file__).parent
# 輸出檔案名稱
OUTPUT_PPT = PROJECT_DIR / "頻審_新展示.pptx"

# 風格配色
COLORS = {
    "primary": RGBColor(30, 60, 110),   # 深藍
    "accent": RGBColor(220, 53, 69),   # 紅
    "text_dark": RGBColor(33, 37, 41), # 深灰
    "bg_light": RGBColor(248, 249, 250) # 淺背景
}

def add_title_slide(prs, title, subtitle):
    """封面幻燈片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景色
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["primary"]
    # 標題文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    # 副標題
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullet_points, accent_color=None):
    """內容幻燈片，使用項目符號"""
    if accent_color is None:
        accent_color = COLORS["accent"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景白色
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg_light"]
    # 標題條
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS["primary"]
    title_bar.line.color.rgb = COLORS["primary"]
    # 標題文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    # 左側色條
    line_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.8), Inches(0.08), Inches(6.7))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = accent_color
    line_shape.line.color.rgb = accent_color
    # 內容文字框
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.8), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["text_dark"]
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    # 1. 封面
    add_title_slide(prs, "頻審專案展示", "功能介紹、技術架構、使用指南")
    # 2. 功能介紹
    add_content_slide(prs, "功能介紹", [
        "✅ 完整的 Discord Bot，支援 AI 對話、經濟系統、遊戲模組",
        "✅ 多角色模式：閒談、數理、語文、程式、家務",
        "✅ 用戶資料持久化，支援金幣、寵物、排行榜等功能",
        "✅ 超過 28 項核心指令，涵蓋管理、互動、娛樂等需求"
    ])
    # 3. 技術架構
    add_content_slide(prs, "技術架構", [
        "🛠️ 後端：Python 3.8+、discord.py >=2.0.0、非同步事件驅動",
        "🤖 AI：Ollama Qwen2.5-1.5B-Instruct，HTTP 非同步呼叫",
        "💾 資料：JSON (user_data.json) 本地持久化",
        "⚙️ 部署：Dockerfile、ngrok 代理，支援雲端執行"
    ], accent_color=COLORS["primary"])
    # 4. 使用指南
    add_content_slide(prs, "使用指南", [
        "1️⃣ 安裝依賴：`pip install -r requirements.txt`",
        "2️⃣ 設定環境變數：編輯 `設定.env`，填入 Discord Bot Token 與 Ollama 端點",
        "3️⃣ 啟動 Bot：`python bot.py` 或使用 Docker `docker compose up -d`",
        "4️⃣ 常用指令：`!help` 查看全部指令，`!role <模式>` 切換角色",
        "5️⃣ 管理員指令：`!ban`、`!kick`、`!shutdown` 等"
    ])
    # 儲存檔案
    prs.save(str(OUTPUT_PPT))
    print(f"✅ PPT 已產生於 {OUTPUT_PPT}")

if __name__ == "__main__":
    create_ppt()
