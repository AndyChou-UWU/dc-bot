"""
生成評審版 PowerPoint - 穩重、專業、簡潔
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ 缺少 python-pptx，正在安裝...")
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

PROJECT_DIR = Path(__file__).parent
PPT_OUTPUT = PROJECT_DIR / "評審版_項目展示.pptx"

# 專業配色
COLORS = {
    'primary': RGBColor(30, 60, 110),      # 深藍
    'accent': RGBColor(220, 53, 69),       # 紅
    'success': RGBColor(40, 167, 69),      # 綠
    'warning': RGBColor(255, 193, 7),      # 黃
    'info': RGBColor(23, 162, 184),        # 青
    'text_dark': RGBColor(33, 37, 41),     # 深灰
    'text_light': RGBColor(108, 117, 125), # 淺灰
    'bg_light': RGBColor(248, 249, 250)    # 淺背景
}

def add_title_slide(prs, title, subtitle):
    """專業標題幻燈片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points, accent_color=None):
    """內容幻燈片"""
    if accent_color is None:
        accent_color = COLORS['accent']
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題背景條
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['primary']
    title_bar.line.color.rgb = COLORS['primary']
    
    # 標題文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 左邊的色條
    line_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.8), Inches(0.08), Inches(6.7))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = accent_color
    line_shape.line.color.rgb = accent_color
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.8), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if isinstance(point, tuple):
            # (縮進級別, 文本)
            level, text = point
            if i == 0 and level == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = text
            p.level = level
        else:
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_feature_grid_slide(prs, title, features):
    """功能網格幻燈片（3 列）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題背景條
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['primary']
    title_bar.line.color.rgb = COLORS['primary']
    
    # 標題文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 左邊的色條
    line_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.8), Inches(0.08), Inches(6.7))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = COLORS['accent']
    line_shape.line.color.rgb = COLORS['accent']
    
    # 網格佈局 (3 列)
    col_width = 2.8
    row_height = 2.0
    start_x = 0.7
    start_y = 1.3
    
    colors = [COLORS['info'], COLORS['success'], COLORS['warning']]
    
    for idx, feature in enumerate(features):
        col = idx % 3
        row = idx // 3
        
        x = start_x + col * (col_width + 0.3)
        y = start_y + row * (row_height + 0.3)
        
        # 方框背景
        box_shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(col_width), Inches(row_height))
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = colors[col]
        box_shape.line.color.rgb = colors[col]
        box_shape.line.width = Pt(2)
        
        # 文字
        text_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(col_width - 0.3), Inches(row_height - 0.3))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.text = feature
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """對比幻燈片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題背景條
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['primary']
    title_bar.line.color.rgb = COLORS['primary']
    
    # 標題文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 左邊的色條
    line_shape = slide.shapes.add_shape(1, Inches(0), Inches(0.8), Inches(0.08), Inches(6.7))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = COLORS['accent']
    line_shape.line.color.rgb = COLORS['accent']
    
    # 左邊盒子
    left_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(4.4), Inches(5.8))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLORS['info']
    left_box.line.color.rgb = COLORS['info']
    left_box.line.width = Pt(2)
    
    left_text = slide.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(4.1), Inches(5.5))
    left_frame = left_text.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(8)
    
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # 右邊盒子
    right_box = slide.shapes.add_shape(1, Inches(5.1), Inches(1.2), Inches(4.4), Inches(5.8))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLORS['success']
    right_box.line.color.rgb = COLORS['success']
    right_box.line.width = Pt(2)
    
    right_text = slide.shapes.add_textbox(Inches(5.25), Inches(1.35), Inches(4.1), Inches(5.5))
    right_frame = right_text.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(8)
    
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(4)
        p.space_after = Pt(4)

def create_review_ppt():
    """創建評審版 PowerPoint"""
    print("🎬 正在生成評審版 PowerPoint...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==================== 第 1 張：封面 ====================
    add_title_slide(prs,
        "Discord Bot 項目展示",
        "subaso-俗北ㄙㄡˊ | v2.0.0")
    
    # ==================== 第 2 張：項目概述 ====================
    add_content_slide(prs, "項目概述", [
        "➜ 功能完整的 Discord Bot，集成 AI 對話、遊戲系統、經濟系統",
        "",
        "➜ 支持 5 種角色模式（閒談、數理、語文、程式、家務）",
        "",
        "➜ 實現用戶數據持久化、聊天記憶上下文管理",
        "",
        "➜ 提供 28 項核心功能，包括遊戲、互動、管理工具"
    ])
    
    # ==================== 第 3 張：技術棧 ====================
    add_comparison_slide(prs, "技術架構",
        "後端框架", [
            "• Discord.py ≥2.0.0",
            "• Ollama AI 引擎",
            "• Python 3.8+",
            "• 非同步事件驅動架構",
            "",
            "優勢：",
            "✓ 低延遲",
            "✓ 可擴展",
            "✓ 穩定可靠"
        ],
        "數據存儲", [
            "• JSON 格式持久化",
            "• user_data.json 用戶檔案",
            "• 本地文件系統",
            "• 實時同步更新",
            "",
            "特點：",
            "✓ 簡單易維護",
            "✓ 無數據庫依賴",
            "✓ 快速查詢"
        ])
    
    # ==================== 第 4 張：核心功能 ====================
    add_feature_grid_slide(prs, "核心功能模塊", [
        "🤖 AI 對話系統",
        "💰 經濟系統",
        "🎮 遊戲系統",
        "👤 用戶管理",
        "🎭 角色切換",
        "📊 數據統計",
        "⚙️ 管理工具",
        "🔧 系統命令",
        "📈 版本管理"
    ])
    
    # ==================== 第 5 張：AI 對話系統 ====================
    add_content_slide(prs, "AI 對話系統", [
        "🤖 集成 Ollama Qwen2.5-1.5B 模型",
        "",
        "💾 上下文記憶：最近 10 條訊息歷史",
        "",
        "🎭 多角色支持：",
        (1, "閒談模式 - 日常對話"),
        (1, "數理模式 - 數學物理計算"),
        (1, "語文模式 - 文學語文輔導"),
        (1, "程式模式 - 代碼解答"),
        (1, "家務模式 - 生活建議"),
        "",
        "⚡ HTTP 異步調用，響應延遲 < 5 秒"
    ])
    
    # ==================== 第 6 張：經濟系統 ====================
    add_content_slide(prs, "經濟與遊戲系統", [
        "💰 初始金幣：1000 幣",
        "",
        "⛏️ 挖礦：50-500 幣 (60 秒冷卻)",
        "🐟 釣魚：40-400 幣 (45 秒冷卻)",
        "🎰 賭博：50% 概率贏或輸 (30 秒冷卻)",
        "🐾 寵物：購買、繁殖、出售 (可獲利)",
        "",
        "✨ 設計特點：",
        (1, "冷卻時間防止刷屏"),
        (1, "隨機性增加可玩性"),
        (1, "經濟平衡，防止無限增長")
    ])
    
    # ==================== 第 7 張：遊戲系統 ====================
    add_content_slide(prs, "互動與遊戲", [
        "🎮 知識競賽系統",
        (1, "Pokemon 猜謎、Trivia 常識、動漫角色、數字遊戲"),
        "",
        "✨ 互動動作系統",
        (1, "!hug !pat !dance !wave 等 8+ 種互動"),
        "",
        "🏆 排行榜系統",
        (1, "金幣排行、遊戲分數、寵物數量排行"),
        "",
        "📊 數據統計",
        (1, "個人成績、遊戲統計、成就解鎖")
    ])
    
    # ==================== 第 8 張：用戶數據管理 ====================
    add_comparison_slide(prs, "數據持久化策略",
        "實時保存", [
            "✓ 每次操作即時保存",
            "✓ 防止數據丟失",
            "✓ 用戶金幣變化立即記錄",
            "✓ 遊戲分數實時更新",
            "",
            "結果：",
            "✅ 用戶信息安全"
        ],
        "結構化存儲", [
            "✓ user_data.json 中央存儲",
            "✓ 按用戶 ID 組織",
            "✓ 包含：金幣、角色、寵物、分數",
            "✓ 易於備份和遷移",
            "",
            "結果：",
            "✅ 易於維護和擴展"
        ])
    
    # ==================== 第 9 張：命令系統架構 ====================
    add_content_slide(prs, "命令系統架構", [
        "📝 統一命令解析器",
        (1, "4 層匹配機制：直接匹配 → 別名 → 忽略大小寫 → 模糊匹配"),
        "",
        "🎯 命令分類",
        (1, "用戶命令（查詢、遊戲）"),
        (1, "管理命令（設置、配置）"),
        (1, "系統命令（幫助、版本）"),
        "",
        "✅ 優勢",
        (1, "靈活的命令識別"),
        (1, "易於擴展新命令"),
        (1, "友善的錯誤提示")
    ])
    
    # ==================== 第 10 張：代碼架構 ====================
    add_comparison_slide(prs, "架構選擇",
        "簡易版本", [
            "📄 bot.py",
            "",
            "✓ 單文件部署",
            "✓ 代碼清晰易懂",
            "✓ 快速原型開發",
            "✓ 易於測試調試",
            "",
            "適用於：",
            "• 學習階段",
            "• 小規模應用"
        ],
        "模塊版本", [
            "🏗️ prism_bot.py + 模塊",
            "",
            "✓ 功能解耦",
            "✓ 易於維護",
            "✓ 代碼重用",
            "✓ 規模可擴展",
            "",
            "適用於：",
            "• 生產環境",
            "• 複雜應用"
        ])
    
    # ==================== 第 11 張：核心流程 ====================
    add_content_slide(prs, "核心執行流程", [
        "1️⃣ 啟動 → 加載配置、連接 Discord API",
        "",
        "2️⃣ 監聽 → 持續接收用戶訊息事件",
        "",
        "3️⃣ 解析 → 識別命令和參數",
        "",
        "4️⃣ 處理 → 執行對應功能邏輯",
        "",
        "5️⃣ 響應 → 返回結果、更新數據",
        "",
        "6️⃣ 保存 → 數據持久化到 user_data.json"
    ])
    
    # ==================== 第 12 張：版本管理 ====================
    add_content_slide(prs, "版本與更新機制", [
        "📦 版本控制：v2.0.0 語義化版本",
        "",
        "📢 更新通知系統",
        (1, "Bot 啟動時檢查版本"),
        (1, "新版本時廣播通知所有頻道"),
        "",
        "📝 版本記錄",
        (1, "user_data.json 中記錄每個用戶的客戶端版本"),
        "",
        "🔄 升級機制",
        (1, "簡單的版本檢查和更新提示")
    ])
    
    # ==================== 第 13 張：安全與穩定性 ====================
    add_content_slide(prs, "安全性與可靠性", [
        "🔐 輸入驗證",
        (1, "命令參數檢查、數據類型驗證"),
        "",
        "⚡ 錯誤處理",
        (1, "Try-except 異常捕獲、友善的錯誤提示"),
        "",
        "📊 日誌系統",
        (1, "按日期分檔，記錄所有操作和錯誤"),
        "",
        "🔄 冷卻機制",
        (1, "防止指令濫用，確保公平性")
    ])
    
    # ==================== 第 14 張：可擴展性 ====================
    add_content_slide(prs, "未來擴展方向", [
        "📱 多平臺支持",
        (1, "支持 Telegram、Slack 等其他平臺"),
        "",
        "🗄️ 數據庫集成",
        (1, "遷移到 MongoDB/PostgreSQL，支持更大規模數據"),
        "",
        "🌐 Web 儀錶板",
        (1, "實時監控 Bot 狀態、用戶統計、配置管理"),
        "",
        "🔌 插件系統",
        (1, "允許社區開發插件擴展功能")
    ])
    
    # ==================== 第 15 張：項目成果 ====================
    add_feature_grid_slide(prs, "項目成果", [
        "✅ 28 項功能",
        "✅ 穩定運行",
        "✅ 用戶數據保護",
        "✅ 模塊化架構",
        "✅ 易於維護",
        "✅ 完整文檔"
    ])
    
    # ==================== 第 16 張：總結 ====================
    add_title_slide(prs,
        "完整、穩定、易擴展",
        "已準備好生產部署")
    
    prs.save(str(PPT_OUTPUT))
    print(f"✅ 評審版 PowerPoint 已生成: {PPT_OUTPUT}")
    print(f"📊 共 {len(prs.slides)} 張幻燈片")
    return True

if __name__ == "__main__":
    try:
        create_review_ppt()
        print("\n🎉 完成！")
    except Exception as e:
        print(f"❌ 生成失敗: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
