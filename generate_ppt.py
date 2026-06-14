"""
生成 Discord Bot 代碼展示 PowerPoint
將整個項目結構與主要代碼片段轉換為演示文稿
"""

import os
import sys
from pathlib import Path
from datetime import datetime

# 檢查並安裝必要的 library
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ 缺少 python-pptx，正在安裝...")
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

# ==================== 設置 ====================
PROJECT_DIR = Path(__file__).parent
PPT_OUTPUT = PROJECT_DIR / "Discord_Bot_代碼展示.pptx"

# 顏色主題
COLOR_TITLE = RGBColor(52, 73, 94)      # 深藍
COLOR_ACCENT = RGBColor(41, 128, 185)   # 亮藍
COLOR_SUCCESS = RGBColor(39, 174, 96)   # 綠色
COLOR_TEXT = RGBColor(44, 62, 80)       # 深灰

def read_file(filepath: Path, max_lines: int = None) -> str:
    """讀取文件內容"""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            lines = f.readlines()
            if max_lines:
                lines = lines[:max_lines]
            return ''.join(lines)
    except:
        return "[無法讀取文件]"

def add_title_slide(prs: Presentation, title: str, subtitle: str):
    """添加標題幻燈片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白佈局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_TITLE
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    for line in subtitle.split('\n'):
        p = subtitle_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0

def add_content_slide(prs: Presentation, title: str, content: str, is_code: bool = False):
    """添加內容幻燈片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題背景
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_ACCENT
    title_shape.line.color.rgb = COLOR_ACCENT
    
    # 標題文字
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6.3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    if is_code:
        # 代碼顯示（使用等寬字體）
        p = text_frame.paragraphs[0]
        p.text = content
        p.font.name = 'Courier New'
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(41, 41, 41)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    else:
        # 普通文本
        for i, line in enumerate(content.split('\n')):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_TEXT
            p.level = 0

def create_presentation():
    """創建完整的 PowerPoint 展示"""
    print("🎬 正在生成 PowerPoint...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==================== 第 1 張：標題 ====================
    add_title_slide(prs, 
        "🤖 Discord Bot 項目",
        "PRISM Bot v2.0\n多功能 AI 聊天機器人\n\n完整代碼與功能展示")
    
    # ==================== 第 2 張：項目概述 ====================
    overview = """📊 項目信息
• 名稱: PRISM Bot（subaso-俗北ㄙㄡˊ）
• 版本: 2.0.0
• 平台: Discord.py
• AI引擎: Ollama (Qwen2.5-1.5B-Instruct)
• 語言: Python 3.8+

🎯 核心功能
✅ AI 對話系統（5個） - 多角色、多語言、連續對話
✅ 經濟系統（6個） - 挖礦、釣魚、狩獵、賭博、寵物、市場
✅ 娛樂遊戲（6個） - Pokemon、Trivia、動漫、數字遊戲
✅ 管理工具（5個） - 日誌、清理、歡迎、配置、自定義命令
✅ 音樂系統（6個）- 可選模組，需 Lavalink

📦 總計 28 個強大功能"""
    
    add_content_slide(prs, "📊 項目概述", overview, is_code=False)
    
    # ==================== 第 3 張：AI 角色系統 ====================
    roles_content = """🎭 5 種 AI 角色

✅ 閒談 - 友善、幽默的聊天夥伴
   使用自然對話風格，可討論任何日常話題

🔢 數理 - 數學和科學專家  
   用清晰的邏輯和步驟解釋問題

📚 語文 - 語言和文學專家
   專長於語法、詞彙、文學分析和寫作技巧

💻 程式 - 專業程式設計師和軟體工程師
   幫助編寫、除錯和解釋程式碼

🏠 家務 - 家務和生活建議專家
   提供關於烹飪、清潔、整理的實用建議

🌍 多語言支持
• 繁體中文、English、日本語、한국어、Español"""
    
    add_content_slide(prs, "🎭 AI 角色與多語言系統", roles_content, is_code=False)
    
    # ==================== 第 4 張：經濟系統 ====================
    economy_content = """💰 經濟系統架構

🪙 虛擬貨幣: 💎 晶幣
   初始余額: 1000 晶幣

⛏️ 挖礦系統
• 鐵礦: 50 晶幣 (60秒冷卻)
• 金礦: 150 晶幣 (60秒冷卻)
• 鑽石: 500 晶幣 (60秒冷卻)

🐟 釣魚系統
• 鯉魚: 40 晶幣 | 金魚: 120 晶幣 | 龍魚: 400 晶幣

🎰 娛樂遊戲
• 賭博 - 50% 贏率
• 老虎機 - 彩票式獎勵

🐾 寵物系統
• 貓 (🐱) | 狗 (🐕) | 龍 (🐉)
• 餵食、交互、心情值系統"""
    
    add_content_slide(prs, "💰 經濟系統詳解", economy_content, is_code=False)
    
    # ==================== 第 5 張：遊戲系統 ====================
    games_content = """🎮 娛樂遊戲模組

🎯 遊戲類型 (6 個)

1️⃣ Pokémon 猜謎
   - 難度: 簡單 (100分) | 中等 (250分) | 困難 (500分)
   - 時限: 30秒 | 20秒 | 10秒

2️⃣ Trivia 知識競賽
   - 多個知識類別
   - 實時計分系統

3️⃣ 動漫角色猜謎
   - 描述猜角色
   - 積分獎勵

4️⃣ 數字猜測遊戲
   - 1-100 的數字範圍
   - 即時反饋

5️⃣ 動漫查詢
   - 元数据查询

6️⃣ 互動動作
   - 擁抱 (🤗) | 拍拍 (👋) | 跳舞 (🕺)

🏆 全局計分系統"""
    
    add_content_slide(prs, "🎮 娛樂遊戲系統", games_content, is_code=False)
    
    # ==================== 第 6 張：指令系統 ====================
    commands_content = """📋 完整指令列表

🤖 AI 對話
!mode [角色] - 切換 AI 角色 (閒談、數理、語文、程式、家務)
!lan [語言] - 切換回答語言
!ask [問題] - 向 AI 提問 (支持連續對話)
!clear - 清除對話歷史

💰 經濟系統
!balance - 查看余額 | !mine - 挖礦 | !fish - 釣魚
!hunt - 狩獵 | !gamble [金額] - 賭博 | !slots [金額] - 老虎機
!pet [list|adopt|info|feed] - 寵物系統

🎮 遊戲
!pokemon - Pokemon猜謎 | !trivia - Trivia競賽
!anime - 動漫角色 | !number - 數字遊戲
!score - 查看遊戲分數 | !guess [答案] - 提交答案

✨ 互動
!hug [@用戶] - 擁抱 | !pat [@用戶] - 拍拍 | !dance - 跳舞

🛠️ 管理 (ADMIN_IDS 使用者)
!admin stats - 查看統計 | !admin update [版本] - 版本通知
!admin backup - 備份數據 | !admin restart - 重啟 Bot"""
    
    add_content_slide(prs, "📋 指令系統", commands_content, is_code=False)
    
    # ==================== 第 7 張：架構 ====================
    architecture = """🏗️ 項目架構

📂 文件結構
bot.py ......................... 單檔版（簡化版）
prism_bot.py ................... 完整版（模組化入口）
prism_config.py ................ 全局配置與常量
modules_economy.py ............. 經濟系統類
modules_games.py ............... 遊戲系統類
modules_admin.py ............... 管理系統類
設定.env ....................... 環境變量配置
requirements.txt ............... 依賴列表
user_data.json ................. 用戶數據存儲

🔌 依賴
discord.py >= 2.0.0 - Discord 機器人框架
httpx - 非同步 HTTP 客戶端 (Ollama API)
python-dotenv - 環境變量管理
logging - 日誌系統

💾 數據持久化
• 用戶個人設定 (user_data.json)
• 經濟系統狀態 (data/economy.json)
• 遊戲歷史紀錄"""
    
    add_content_slide(prs, "🏗️ 項目架構", architecture, is_code=False)
    
    # ==================== 第 8 張：bot.py 核心代碼 ====================
    bot_code = read_file(PROJECT_DIR / "bot.py", max_lines=30)
    add_content_slide(prs, "💻 bot.py - 核心入口 (前 30 行)", bot_code, is_code=True)
    
    # ==================== 第 9 張：prism_config.py ====================
    config_code = read_file(PROJECT_DIR / "prism_config.py", max_lines=35)
    add_content_slide(prs, "⚙️ prism_config.py - 配置文件 (前 35 行)", config_code, is_code=True)
    
    # ==================== 第 10 張：!mode 指令邏輯 ====================
    mode_snippet = """# !mode 指令 - 角色切換邏輯 (bot.py)

if message.content.startswith("!mode"):
    args = message.content[5:].strip().split()
    if not args:
        available = ", ".join(f"**{k}**" for k in PERSONALITIES.keys())
        await message.channel.send(f"目前角色: **{current_personality}**\\n可用角色: {available}")
        return

    raw_input = args[0].strip()
    normalized = raw_input

    # 直接匹配 → 別名映射 → 大小寫容錯 → 模糊匹配
    if normalized not in PERSONALITIES:
        alias_map = {"chat": "閒談", "math": "數理", ...}
        lowered = raw_input.lower()
        if lowered in alias_map:
            normalized = alias_map[lowered]
        else:
            matches = [k for k in PERSONALITIES.keys() if k.lower() == lowered]
            if matches:
                normalized = matches[0]

    if normalized in PERSONALITIES:
        user_personalities[user_id] = normalized
        save_all_user_data()
        await message.channel.send(f"✅ 已切換到 **{normalized}** 模式")"""
    
    add_content_slide(prs, "🎭 角色切換邏輯", mode_snippet, is_code=True)
    
    # ==================== 第 11 張：AI 請求流程 ====================
    ai_snippet = """# AI 請求處理流程 (bot.py)

async def handle_ai_request(message, question, personality):
    user_id = message.author.id
    if user_id not in user_conversations:
        user_conversations[user_id] = []

    language = user_languages.get(user_id, "chinese")
    system_prompt = PERSONALITIES[personality]["system"] + " " + LANGUAGE_PROMPTS[language]
    
    # 構建消息序列
    api_messages = [{"role": "system", "content": system_prompt}]
    api_messages.extend(user_conversations[user_id])
    api_messages.append({"role": "user", "content": question})

    # 調用 Ollama API
    async with httpx.AsyncClient(timeout=None) as client_http:
        payload = {
            "model": OLLAMA_MODEL,
            "messages": api_messages,
            "stream": False,
            "options": {"num_ctx": 4096}
        }
        response = await client_http.post(f"{OLLAMA_BASE_URL}/api/chat", json=payload)
        result = response.json()
        answer = result.get("message", {}).get("content", "⚠️ AI 未返回內容")

    # 儲存對話歷史
    user_conversations[user_id].append({"role": "user", "content": question})
    user_conversations[user_id].append({"role": "assistant", "content": answer})
    await message.channel.send(f"**【{personality}】**\\n{answer}\""""
    
    add_content_slide(prs, "🤖 AI 請求處理", ai_snippet, is_code=True)
    
    # ==================== 第 12 張：版本管理與廣播 ====================
    version_snippet = """# 版本管理與廣播系統 (bot.py)

async def broadcast_update_notification(new_version):
    '''向所有用戶廣播版本更新通知'''
    global BOT_VERSION, last_notified_version
    BOT_VERSION = new_version
    last_notified_version = new_version
    save_all_user_data()
    
    update_message = f"🎉 **系統更新通知** 🎉\\n✅ 已更新成 **{new_version}**\\n新版本已推送，感謝使用！"
    
    success_count = 0
    for user_id, dm_channel in list(user_dms.items()):
        try:
            await dm_channel.send(update_message)
            success_count += 1
            await asyncio.sleep(0.5)  # 避免速率限制
        except:
            pass
    
    logger.info(f"更新通知已發送: 成功 {success_count}")

# 管理員指令
!admin update 2.0.1  # 向所有用戶廣播版本更新"""
    
    add_content_slide(prs, "📡 版本管理與廣播系統", version_snippet, is_code=True)
    
    # ==================== 第 13 張：數據持久化 ====================
    data_snippet = """# 用戶數據持久化系統 (bot.py)

def load_all_user_data():
    '''從 JSON 加載所有用戶數據'''
    global user_personalities, user_languages
    data = load_user_data()
    if "personalities" in data:
        user_personalities = {int(k): v for k, v in data["personalities"].items()}
    if "languages" in data:
        user_languages = {int(k): v for k, v in data["languages"].items()}
    logger.info(f"加載了 {len(user_personalities)} 個用戶的數據")

def save_all_user_data():
    '''保存所有用戶數據到 JSON'''
    data = {
        "personalities": {str(k): v for k, v in user_personalities.items()},
        "languages": {str(k): v for k, v in user_languages.items()},
        "last_save": datetime.now().isoformat(),
        "bot_version": BOT_VERSION,
        "last_notified_version": last_notified_version
    }
    save_user_data(data)

# 數據文件位置: user_data.json"""
    
    add_content_slide(prs, "💾 數據持久化系統", data_snippet, is_code=True)
    
    # ==================== 第 14 張：特性與優勢 ====================
    features = """✨ 核心特性

🔄 連續對話
   • 保持上下文歷史（最近 10 次交互）
   • 自動清理超出限制的歷史記錄

🌐 多語言多角色
   • 5 種角色各有專屬系統提示
   • 支持 5 種語言即時切換
   • 用戶個人設定獨立存儲

💾 完整的數據持久化
   • 用戶配置（角色、語言）自動保存
   • JSON 格式便於檢查和調試
   • 版本管理與更新通知

⚡ 非同步架構
   • 基於 asyncio + discord.py 2.0+
   • 支持並發消息處理
   • 高效的 HTTP 請求（httpx）

🛡️ 容錯機制
   • 角色切換支持大小寫、別名、模糊匹配
   • 指令解析的容錯設計
   • 完善的異常處理與日誌記錄"""
    
    add_content_slide(prs, "✨ 核心特性與優勢", features, is_code=False)
    
    # ==================== 第 15 張：部署與使用 ====================
    deployment = """🚀 部署與使用步驟

1️⃣ 環境設置
   python -m venv venv
   source venv/bin/activate  # Linux/Mac
   venv\\Scripts\\activate     # Windows
   pip install -r requirements.txt

2️⃣ 配置環境變量 (設定.env)
   DISCORD_TOKEN=your_token_here
   OLLAMA_BASE_URL=http://localhost:11434
   OLLAMA_MODEL=Qwen2.5-1.5B-Instruct
   ADMIN_IDS=123456789,987654321

3️⃣ 啟動 Ollama (另一個終端)
   ollama serve

4️⃣ 運行 Bot
   python bot.py        # 簡化版
   python prism_bot.py   # 完整版（推薦）

5️⃣ 使用
   • 添加 Bot 到 Discord 伺服器
   • 在私信中使用 !help 查看所有命令
   • 開始聊天！

📝 文件位置
logs/              - 日誌文件
user_data.json     - 用戶數據
設定.env           - 環境配置"""
    
    add_content_slide(prs, "🚀 部署與使用", deployment, is_code=False)
    
    # ==================== 第 16 張：總結 ====================
    summary = """📌 項目總結

🎯 目標達成
✅ 完整的 Discord Bot 框架
✅ 28 個功能的完整實現
✅ 生產級別的代碼品質
✅ 完善的用戶數據管理
✅ 靈活的 AI 角色系統

💪 技術亮點
✓ 非同步編程最佳實踐
✓ 模組化與可擴展設計
✓ 多層次的容錯機制
✓ 完整的日誌系統
✓ 版本控制與廣播通知

🔮 未來可能的擴展
• 集成更多 Ollama 模型
• 添加 Web Dashboard
• 支持語音交互
• 機器學習模型優化
• 多服務器支持

📞 開發者信息
名稱: subaso-俗北ㄙㄡˊ
版本: 2.0.0
更新日期: 2026-05-10
GitHub: https://github.com/AndyChou-UWU/dc-bot"""
    
    add_content_slide(prs, "📌 項目總結", summary, is_code=False)
    
    # ==================== 保存 ====================
    prs.save(str(PPT_OUTPUT))
    print(f"✅ PowerPoint 已生成: {PPT_OUTPUT}")
    print(f"📊 共 {len(prs.slides)} 張幻燈片")
    return True

if __name__ == "__main__":
    try:
        create_presentation()
        print("\n🎉 完成！")
    except Exception as e:
        print(f"❌ 生成失敗: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
